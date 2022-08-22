"""
Parse AFWERX Proposals for the following info an aggregate it in one .csv file:
1. Firm Certificate Questions
2. Proposal Certification Questions
3. DAF customer and end-user info
4. TPOC info
5. Various budget line items
"""
import argparse
import os
import sys
from utils import is_directory, is_filename

# TODO Update with mandatory sections, then check for their presence
sections = [
    "Method",
    "Approach"
]

key_phrases = [
    "DAF Customer",
    "DAF End-User",
    "Digitally signed by",
    "TPOC:",
    "TPOCs:",
    "TPOCS:",
    "Technical Point of Contact"
]

# ==============================================================================
# Reference:
# https://stackoverflow.com/questions/15008758/
def str2bool(this_string):
    """
    Validates that an argparse argument is a boolean value.
    """
    if isinstance(this_string, bool):
        return this_string
    if this_string.lower() in ('yes', 'true', 't', 'y', '1'):
        return True
    elif this_string.lower() in ('no', 'false', 'f', 'n', '0'):
        return False
    else:
        raise argparse.ArgumentTypeError('Boolean value expected.')
# ==============================================================================
def lower_str(x):
    """
    Converts a provided argparse option string to a lowercase string.
    """
    return x.lower()

# ==============================================================================
def process_ppt(file_name):
    """
    Prints the title of every slide
    """
    # Lazily load the pptx module - this will only execute once if needed at all
    from pptx import Presentation
    prs = Presentation(file_name)

    for slide_index, slide in enumerate(prs.slides):
        try:
            title = slide.shapes.title.text
        except AttributeError:
            title = slide.shapes[0].text
        print(f"Slide #{slide_index}: {title=}")
# ==============================================================================
def process_pdf_page_titles(file_name):
    """
    Prints the title of every page (intended for slides in pdf format)
    """
    text = ''
    with fitz.open(file_name) as doc:
        for page_index, page in enumerate(doc):
            logger.debug(f"{page_index}".center(80,"-"))
            text = page.get_text().split('\n')[0]
            print(f"Slide #{page_index}: {title=}")
# ==============================================================================
def parse_firm_certificate(seg_i, single_text, text_segs, firm_cert_questions):
    """
    Given a segment index, single line of text, list of text segments,
    and firm certificate questions
    (provided as a dict mapping question number to (partial) question text),
    check if this single line of text matches any of the firm certification
    questions.  If it matches any of them, parse the corresponding answer and
    return it as a {question number:question answer} dict.

    If no match, return an emtpy dictionary: {}
    """
    result = {}
    answer = ""
    for value, key in firm_cert_questions.items():
        if not key in single_text:
            continue

        if value == 6:
            if key in single_text:
                text = text_segs[seg_i+1]
                result[value] = text.strip()

        elif value in [7,8]:
            if key in single_text:
                text = text_segs[seg_i+2]
                result[value] = text.strip()

        elif value in [10,11]:
            result[value] = ""
            text = text_segs[seg_i]
            for index in range(10):
                if '[X]' in text:
                    result[value] += text.strip() + '\n'
                text = text_segs[seg_i+index]

        elif value == 16:
            if key in single_text:
                index = -2
                while index < 8:
                    answer = text_segs[seg_i+index].strip()
                    # TODO clean this logic up
                    if some_digits.search(answer) and not answer.endswith("pdf"):
                        result = {value:answer}
                        break
                    index += 1
        else:
            index = 1
            while True:
                answer = text_segs[seg_i+index].strip()
                if answer:
                    result = {value:answer}
                    break
                index += 1
            break

    return result
# ==============================================================================
def parse_proposal_certification(   seg_i,
                                    single_text,
                                    text_segs,
                                    prop_cert_questions,
                                    prop_type,
                                    last_answer):
    """
    Given a segment index, single line of text, list of text segments,
    and proposal certification questions
    (provided as a dict mapping question number to (partial) question text),
    check if this single line of text matches any of the firm certification
    questions.  If it matches any of them, parse the corresponding answer and
    return it as a {question number:question answer} dict.

    If no match, return an empty dictionary: {}
    """
    result = {}
    answer = ""

    for key, value in prop_cert_questions.items():
        if value not in single_text or value == "N/A":
            continue

        if single_text.split()[-1] in ["YES", "NO"]:
            answer = single_text.split()[-1]

        elif prop_type == "SBIR":
            #
            if key in [3,4]:
                answer = text_segs[seg_i+2].strip()
                if not answer:
                    answer = text_segs[seg_i+3].strip()

            elif 6 == key:
                if value in single_text:
                    for text in text_segs[seg_i+1:seg_i+3]:
                        answer = text.strip()
                        if answer.lower() in ["yes", "no"]:
                            break
            else:
                index = 1
                while index < 8:
                    answer = text_segs[seg_i+index].strip()
                    if answer:
                        break
                    index += 1

        elif prop_type == "STTR":
            if key == 4:
                answer = text_segs[seg_i+2].strip()
                if not answer:
                    answer = text_segs[seg_i+3].strip()

            elif 6 == key:
                if value in single_text:
                    for text in text_segs[seg_i+1:seg_i+3]:
                        answer = text.strip()
                        if answer.lower() in ["yes", "no"]:
                            break
            else:
                index = 1
                while index < 8:
                    answer = text_segs[seg_i+index].strip()
                    if answer:
                        break
                    index += 1
        if len(answer.split()) != 1:
            logger.info(f"Answer to Prop Q #{key} of '{answer}' seems wrong, patching over it with last YES/NO text block found:'{last_answer}'")
            answer = last_answer
        if key < 16 and answer.lower() not in ["no", "yes"]:
            logger.error(f"Parsed answer to Proposal Q #{key} not YES/NO. Skipping")
        else:
            result = {key:answer}
            logger.debug(str(result))
    return result
# ==============================================================================

def parse_safety(seg_i, single_text, text_segs):
    """
    Parse the "Safety Related Deliverables" section, returning all paragraphs
    under that heading.
    """
    start_number = 0
    result = ""
    low_start = 0
    disclaimer = "Use or disclosure of data contained on this page is subject to the restriction on the first page of this volume"
    # Return immediately if that header isn't present
    if not safety_heading.search(single_text.lower()):
        return {}

    # Header is present, search nearby for an associated number, e.g.
    # 2.7 Safety Related Deliverables
    for i in range(-2,1):
        try:
            start_number = re.match("(\d(\.\d)+)", text_segs[seg_i+i].split()[0]) #float(text_segs[seg_i+i].split()[0].rstrip('.'))
            if start_number:
                start_number = start_number.group()
                low_start = i
                break
        except (ValueError, IndexError) as e:
            pass
            print(f"Exception {e}")

    # Low start <0 means the number was found in a
    # text block BEFORE the key word
    if low_start < 0:
        for i in range(low_start,0,1):
            result += text_segs[seg_i+i] + ' '
        #result += '\n'

    # Keep up to 100 lines that follow, searching for the next section, e.g.
    #  the next number, e.g.
    # 2.8 Other Deliverables
    # 3. Next section
    # Stop when this next section is found
    # TODO How many lines to get?  4? 10?
    print(f"{start_number=}")
    for i in range(20):
        try:
            if bool(re.search("\d\.|^\d$|Commercialization Strategy", text_segs[seg_i+i].strip())):
                number_str = text_segs[seg_i+i].split()[0]#.rstrip('.')
                #new_number = float(number_str)
                new_number = number_str
                print(f"{new_number=}")
                """
                delta = math.floor(new_number) - math.floor(start_number)
                if delta >= 1:
                    logging.debug(f"Safety, found new heading: {new_number=} {start_number=} on line '{text_segs[seg_i+i]}'")
                    break
                """
                if not number_str.startswith(str(start_number)):
                    print(f"{number_str} doesn't start with {start_number}; breaking")
                    break
        except ValueError as e:
            print(f"Error: {e=}")
            pass
        if disclaimer != text_segs[seg_i+i]:
            result += text_segs[seg_i+i] + '\n'
            logger.debug(text_segs[seg_i+i])

    if result:
        print(f"Safety: {result=}")
        return {"Safety-Related Deliverables": result}

    return {}
# ==============================================================================
def parse_type_phase(text):
    phase = ""
    p_type = ""
    if "Phase III" in text:
        phase = "III"
    elif "Phase II" in text:
        phase = "II"
    elif "Phase I" in text:
        phase = "I"

    if "SBIR" in text:
        p_type = "SBIR"
    elif "STTR" in text:
        p_type = "STTR"
    return p_type, phase
# ==============================================================================
def parse_questions(file_name):
    """
    Parse all relevant questions from the all_forms files, e.g.
    F2D-1234_All_forms_proposal_package.pdf
    """
    #print(f"Parsing all forms: {file_name}")

    sbir_phase_I_prop_cert_question_1 = "Written Approval:"
    sbir_phase_II_prop_cert_question_1 = "officer:"

    sbir_prop_cert_questions = {
        1: "officer:",
        2: "705?",
        3: "During the performance of the contract, the research/research and development will be performed",
        4: "offerors facilities by the offerors employees except as otherwise indicated in the technical",
        5: "or equipment?",
        6: "control regulations",
        7: "There will be ITAR/EAR data in this work and/or deliverables.",
        8: "components?",
        9: "proposals listed above",
        10: "another Federal agency",
        11: "disclosure restriction?",
        12: "DNA of the solicitation",
        13: "without evaluation",
        14: "subcontractors proposed",
        15: "22 CFR 120.16",
        16: "will be on the project?",
        17: "Is the principal investigator socially/economically disadvantaged?",
        18: "Economic Development Organizations?",
        19: "N/A",
        20: "N/A",
        21: "N/A",
        22: "N/A"
    }

    sttr_prop_cert_questions = {
        1: "Is the institution a Research Institution?",
        2: "the research institute as described in 13 C.F.R",
        3: "as defined by 13 C.F.R",
        4: "research/research and development will be performed in the",
        5: "offerors employees except as otherwise indicated in the technical proposal.",
        6: "Do you plan to use Federal facilities, laboratories, or equipment?",
        7: "The offeror understands and shall comply with export control regulations",
        8: "There will be ITAR/EAR data in this work and/or deliverables",
        9: "components",
        10: "Contract been awarded for any of the proposals listed above",
        11: "under this award is subsequently funded by another Federal agency",
        12: "disclosure restriction",
        13: "Recombinant DNA of the solicitation:",
        14: "may be cause for rejection of the proposal submission without evaluation.",
        15: "Are teaming partners or subcontractors proposed?",
        16: "CFR 120.16 for work under the proposed effort?",
        17: "Percentage of the principal investigators total time will be on the project:",
        18: "Firm Percentage of Work",
        19: "Research Institute Percentage of Work",
        20: "Other Subcontractor Percentage of Work",
        21: "Is the principal investigator socially/economically disadvantaged?",
        22: "contact information to Economic Development Organizations?"
    }

    firm_cert_questions = {
        1: "requirements set forth in 13 C.F.R.",
        2: "requirements are U.S. citizens or permanent resident aliens in the United States.",
        3: "It has no more than 500 employees, including the employees of its affiliates.",
        4: "Number of employees including all affiliates (average for preceding 12 months)",
        5: "It has met the performance benchmarks as listed by the SBA on their website as eligible to participate",
        6: "funds or private equity",
        7: "It has more than 50% owned by a single Venture Capital Owned Company (VCOC), hedge fund, or private equity",
        8: "It has more than 50% owned by multiple business concerns that are VOCs, hedge funds, or private equity",
        9: "Firms PI, CO, or owner, a faculty member or student of an institution of higher education",
        10: "The offeror qualifies as a:",
        11: "Race of the offeror:",
        12: "Ethnicity of the offeror",
        13: "responsible for collecting the tax liability:",
        14: "involving federal funds",
        15: "for a fraud-related violation involving federal funds:",
        16: "Supporting Documentation:",
        17: "firm owned or managed by a corporate entity?",
        18: "Is your firm affiliated as set forth in 13 CFR"
    }

    duration = "Proposed Base Duration (in months)"
    prop_cert_questions = None

    text_segs = []
    safety_info_found = 0
    result = {}
    answer = ""
    type_found = False
    prop_type = False
    last_answer = None
    with fitz.open(file_name) as doc:
        for page_i, page in enumerate(doc):
            logger.debug(f"{page_i}".center(80,"-"))
            text_segs += page.get_text().split('\n')

        #for seg_i, (page_i, single_text) in enumerate(text_segs):
        #    print(f"Page #{page_i}, text #: {seg_i}: '{single_text}'")

        # Iterate over every text block in the PDF
        for seg_i, single_text in enumerate(text_segs):
            if not single_text:
                continue

            if not type_found and "Phase" in single_text and "Proposal" in single_text:
                type_found = True
                prop_type, prop_phase = parse_type_phase(single_text)
                if prop_type == "SBIR":
                    prop_cert_questions = sbir_prop_cert_questions
                elif prop_type == "STTR":
                    prop_cert_questions = sttr_prop_cert_questions
                this_answer = {"Type":prop_type, "Phase": prop_phase}
                logger.debug(f"{single_text=} parsed to {this_answer=}")
                if this_answer["Type"] == "SBIR" and this_answer["Phase"] == "I":
                    logger.debug("Updating SBIR Prop Q1")
                    sbir_prop_cert_questions[1] = sbir_phase_I_prop_cert_question_1
                result.update(this_answer)

            if not prop_type:
                continue

            single_text = single_text.strip()
            if len(single_text.split()) == 1:
                last_answer = single_text
                #continue

            # Remove entries from the firm_cert_questions list once found
            if firm_cert_questions:
                firm_cert_info = parse_firm_certificate(seg_i,
                                                        single_text,
                                                        text_segs,
                                                        firm_cert_questions)
                if firm_cert_info:
                    value, answer = firm_cert_info.popitem()
                    this_answer = {f"Firm Certification Q{value}":answer}
                    firm_cert_questions.pop(value)
                    result.update(this_answer)
                    continue

            # Remove entries from the prop_cert_questions list once found
            if prop_cert_questions:
                prop_cert_info = parse_proposal_certification(  seg_i,
                                                                single_text,
                                                                text_segs,
                                                                prop_cert_questions,
                                                                prop_type,
                                                                last_answer)
                if prop_cert_info:
                    value, answer = prop_cert_info.popitem()
                    this_answer = {f"Proposal Certification Q{value}":answer}
                    prop_cert_questions.pop(value)
                    result.update(this_answer)
                    if result['Type'] == "SBIR" and 19 in prop_cert_questions:
                        for key in range(19,23):
                            result[f"Proposal Certification Q{key}"] = prop_cert_questions.pop(key)
                    continue

            # This safety info appears twice, once in the table of contents
            #and once in the body
            if safety_info_found < 2:
                safety_info = parse_safety(seg_i, single_text, text_segs)
                if safety_info:
                    safety_info_found +=1
                    result.update(safety_info)
                    continue

            if duration and duration.lower() in single_text.lower():
                result["Duration (Mo.)"] = single_text.split()[-1].strip()
                duration = False

    #for prop_cert_question in prop_cert_questions:
    logger.debug(f"{result=}")
    if prop_cert_questions:
        pprint.pprint(prop_cert_questions)
        result["Missing Proposal Certificate questions"] = list(prop_cert_questions.keys())
    if firm_cert_questions:
        print(f"Couldn't find Proposal Certificate questions:")
        pprint.pprint(firm_cert_questions)
        result["Missing Firm Certificate Questions"] = list(firm_cert_questions)
    return result
# ==============================================================================
def parse_budget(file_name,
                 max_value,
                 keyphrase="Total Dollar Amount for this Proposal"):
    """
    Parse all relevant fields from budget (see lists below)

    Why make a separate function here when I could parse it from the
    all_forms file? Because the budget file is much smaller, so it *should be*
    faster to search for the budget info.  I think.
    # TODO Verify the above with imperical testing.
    """
    #print(f"Parsing budget: {file_name}")
    summed_costs = defaultdict(float)
    sumable_headings =   [
        "Total Direct Travel Costs (TDT)",
        "Total Direct Material Costs (TDM)",
        "Total Subcontractor Costs (TSC)",
        "Total Direct Supplies Costs (TDS)",
        "Total Direct Equipment Costs (TDE)",
        "Total Other Direct Costs (TODC)",
        "Total Direct Labor (TDL)",
    ]
    total_heading = "Total Dollar Amount for this Proposal"
    text_segs = []
    total_proposal_cost = 0
    unique_costs = set()
    result = {}
    with fitz.open(file_name ) as doc:
        for page in doc:
            text_segs += page.get_text().split('\n')

        for seg_i, single_text in enumerate(text_segs):
            logger.debug(single_text)
            logger.debug(f"{keyphrase} {type(keyphrase)}")

            for heading in sumable_headings:
                if heading.lower() in single_text.lower():
                    cost_str = text_segs[seg_i+1]
                    cost_float = float(cost_str.lstrip('$').replace(",",""))
                    if not cost_float or cost_float not in unique_costs:
                        unique_costs.add(cost_float)
                        summed_costs[heading] += cost_float

            if total_heading and total_heading.lower() in single_text.lower():
                logger.debug(single_text)
                budget_str = text_segs[seg_i+1]
                logger.debug(budget_str)
                total_proposal_cost = float(budget_str.lstrip('$').replace(",",""))
                result["Total"] = total_proposal_cost
                logger.debug(result)
                if total_proposal_cost > max_value:
                    print(f"WARNING! Proposed budget exceeds ${max_value}!")
                total_heading = False
                continue

    #pprint.pprint(summed_costs)
    result.update(summed_costs)
    logger.debug(result)
    return result
# ==============================================================================
def get_total_budget(file_name,
                     max_value=1250000,
                     keyphrase="Total Dollar Amount for this Proposal"):
    """
    Print total budget shown in budget document
    """
    budget_float = -1.0
    threshold = "$"+str(max_value/1000000)+"M"
    with fitz.open(file_name ) as doc:
        for page in doc:
            logger.debug(f"{page_count}".center(80,"-"))
            text_segs = page.get_text().split('\n')
            for seg_i, single_text in enumerate(text_segs):
                logger.debug(keyphrase, type(keyphrase))
                if keyphrase.lower() not in single_text.lower():
                    continue

                logger.debug(single_text)
                budget_str = text_segs[seg_i+1]
                logger.debug(budget_str)
                budget_float = float(budget_str.lstrip('$').replace(",",""))
                if budget_float > max_value:
                    logger.error(f"WARNING! Proposed budget exceeds ${threshold}!")
                break
    return {"Total Proposal Value" :budget_float}
# ==============================================================================
def process_pdf_sigs_fitz(file_name):
    """
    Collect POC info, returning immediately once all 3 POCs found
    """
    got_text = False
    result = defaultdict(str)
    headers = [
        "Primary End-User Organization",
        "Primary Customer Organization",
        "Technical Points of Contact (TPOCs)"
    ]
    with fitz.open(file_name) as doc:

        # Iterate over every page in the doc
        for pi, page in enumerate(doc):
            text_segs = page.get_text().split('\n')
            #text_segs = [text.strip() for text in text_segs if text.strip()]
            if not text_segs:
                continue

            got_text = True
            logger.debug(text_segs)
            # Iterate over every text field
            for seg_i, single_text in enumerate(text_segs):
                logger.debug(f"Page #{pi}, text #: {seg_i}: '{single_text}'")
                for key_phrase in headers:
                    if key_phrase not in single_text:
                        continue

                    count = 0
                    index = 0
                    # Sometimes there are blank newlines
                    # Skip over them, grabbing the next two lines with text
                    while index < 10:
                        try:
                            x = text_segs[seg_i+index]
                        except IndexError as e:
                            break
                        if x.strip():
                            count += 1
                            result[key_phrase] += x + '\n'
                            logger.debug(x)

                        index +=1
                        if x.rstrip().endswith(","):
                            continue
                        if count > 2 or index > 10:
                            break

                    headers.remove(key_phrase)
                    if not headers:
                        return result
                    continue
    return result
# ==============================================================================
def ocr_cleanup(open_file_handle, open_file_name, files_to_remove):
    """
    Clean up artifacts of OCR - temp files & open file handles
    """
    print(f"Removing OCR artifacts related to {open_file_name}")
    for file_name_to_remove in files_to_remove:
        os.remove(file_name_to_remove)
    if open_file_handle:
        open_file_handle.close()
    if open_file_name:
        os.remove(open_file_name)
# ==============================================================================
def ocr_pdf(file_name):
    """
    #TODO Don't parse the digital signature section, rather,
    parse the paragraph headings instead as in process_pdf_sigs_fitz()

    Use optical character recognition to parse scanned PDFs
    https://www.geeksforgeeks.org/python-reading-contents-of-pdf-using-ocr-optical-character-recognition/
    """
    print(f"OCR'ing {file_name}.  This could take a minute.")
    # Counter to store images of each page of PDF to image
    image_counter = 1
    files_to_remove = []
    pages = convert_from_path(file_name, 500)
    # Iterate through all the pages stored above
    for page in pages:
        filename = "page_"+str(image_counter)+".jpg"
        # Save the image of the page in system
        page.save(filename, 'JPEG')
        files_to_remove.append(filename)
        image_counter += 1

    #Part #2 - Recognizing text from the images using OCR
    # Creating a text file to write the output
    outfile = "out_text.txt"

    # Open the file in append mode so that
    # All contents of all images are added to the same file
    f = open(outfile, "a")

    # Iterate from 1 to total number of pages
    for filename in files_to_remove:
        # Recognize the text as string in image using pytesserct
        text = str(((pytesseract.image_to_string(Image.open(filename)))))
        # Finally, write the processed text to the file.
        f.write(text)
        print(f"Page {filename}")

        # The recognized text is stored in variable text
        # text = text.replace('-\n', '')
        text_segs = text.split("\n")
        text_segs = [x.strip() for x in text_segs if x]
        # Iterate over every text field
        for seg_i,_  in enumerate(text_segs):
            single_text = text_segs[seg_i]
            for key_phrase in key_phrases:
                if key_phrase in single_text:
                    print(f"Segment {seg_i}".center(80, "*"))
                    for x in text_segs[seg_i-2:seg_i+5]:
                        print(x)
                    print("*"*80)
                    logger.debug(single_text.strip())
                    #ocr_cleanup(f, outfile, files_to_remove)
                    #return[0]
            logger.debug(text)
    ocr_cleanup(f, outfile, files_to_remove)
# ==============================================================================
def parse_file(file_name, prop_number, ocr_flag):
    """
    Called by main(): this is the top level function for processing any file.
    Having one top-level function in this fashion allows for easy
    parallelization with a multiprocessing Pool later, once all other
    optimizations are complete.
    # TODO Call this in parallel.
    """
    logger.debug("-"*80)
    file_info = {}
    temp_info = {}#{"Phase":'?'}
    poc_info = {}
    file_extension = file_name.split(".")[-1]
    short_name = file_name.split("\\")[-1]
    # Process PowerPoint files
    if file_extension in ppt_extensions:
        process_ppt(file_name)
        #total_files +=1

    # All others are PDFs
    else:
        #process_pdf_page_titles(file_name)
        #if any(keyword in file_name.lower() for keyword in args.questions_file):
        if args.questions_file.lower() in file_name.lower():
            logger.info(f"Parsing {short_name} for questions")
            temp_info.update(parse_questions(file_name))
            file_info.update(temp_info)
            if not temp_info:
                if not file_info and ocr_flag:
                    ocr_pdf(file_name)
                else:
                    print(f"Can't parse {short_name}; Consider enabling OCR with -o True")
        #if "budget" in file_name:

        if args.budget_file.lower() in file_name.lower():
            logger.info(f"Parsing {short_name} for budget info")
            temp_info.update(parse_budget(file_name, args.max_value))
            file_info.update(temp_info)
            #file_info.update(get_total_budget(file_name))
            if not temp_info:
                if not file_info and ocr_flag:
                    ocr_pdf(file_name)
                else:
                    print(f"Can't parse {short_name}; Consider enabling OCR with -o True")

        # Try to get signatures and TPOC data from this PDF
        # If all 3 POCs haven't yet been found, search this file for them
        if args.sig_file.lower() in file_name.lower():
            logger.info(f"Parsing {short_name} for signatures")
            poc_info = process_pdf_sigs_fitz(file_name)
            if poc_info:
                print(f"Found POC info in {short_name}")
                file_info.update(poc_info)

        # Count number of pages/slides in the proposal document
        # If all 3 POCs haven't yet been found, search this file for them
        if args.vol2_file.lower() in file_name.lower():
            logger.info(f"Counting slides/pages in {short_name}")
            page_i = 0
            with fitz.open(file_name) as doc:
                for page_i, page in enumerate(doc):
                    pass
            file_info["Page Count"] = page_i

    for k in file_info.keys():
        try:
            file_info[k] = file_info[k].strip()
        except AttributeError:
            pass
    return file_info
# ==============================================================================
#https://www.tutorialspoint.com/
# How-to-correctly-sort-a-string-with-a-number-inside-in-Python
def atoi(text):
    """
    Helper function for sorting strings with numbers within them
    """
    return int(text) if text.isdigit() else text
#==============================================================================
#https://www.tutorialspoint.com/
# How-to-correctly-sort-a-string-with-a-number-inside-in-Python
def natural_keys(text):
    """
    Helper function #2 for sorting strings with numbers within them
    """
    return [atoi(c) for c in re.split('(\d+)', text)]
#==============================================================================
def done_gathering_info(info):
    # TODO
    pass
#==============================================================================
def main():
    """
    Create a list of files as provided by -f and -d flags.
    Will recursively traverse all -d directories keeping files whose names
    include all terms specified by the -k options.

    Write results to args.out
    """
    target_files = set()
    dirs = []
    prop_status = defaultdict(str)
    all_info = defaultdict(dict)
    poc_headers = [
        "Primary End-User Organization",
        "Primary Customer Organization",
        "Technical Points of Contact (TPOCs)"
    ]

    if not args.directory:
        args.directory = set()
    if not args.file:
        args.file = set()
    for directory in args.directory:
        dirs.append(directory)

    for file_name in args.file:
        file_extension = file_name.split(".")[-1]
        if file_extension in valid_extensions:
            target_files.add(file_name)
        else:
            print(f"Skipping {file_name} with extension {file_extension}")

    # Count the total number of files to be parsed by recursively walking
    # all provided directories.
    for source_dir in dirs:
        for (root, _, files) in os.walk(source_dir):
            for file_name in files:
                file_extension = file_name.split(".")[-1]
                if ((args.keyword and all(x in file_name.lower() for x in args.keyword)) or \
                    not args.keyword) and \
                    file_extension.lower() in valid_extensions:
                    target_files.add((root+'\\'+file_name, os.path.getsize(root+'/'+file_name)))

    print(f"Parsing {len(target_files)} files. This could take a few seconds.")
    # Reset the counter.  It will be incremented as each file is parsed.
    #total_files = 0

    # Here is the serial (non-parallel) approach.  Slow, but it works.
    start = time.time()
    # TODO Parallelize
    target_files = sorted(target_files, key=lambda x:x[1], reverse=True)

    for file_name, _ in sorted(target_files):
        prop_number = file_name.split("\\")[-1].split("_")[0]#re.search(four_digits, file_name)
        if not prop_number or prop_status[prop_number]:
            continue
        #prop_number = prop_number.group(1)
        logger.debug(f"Proposal: {prop_number}")
        temp_info = parse_file(file_name, prop_number, args.ocr)
        if done_gathering_info(temp_info):
            prop_status[prop_number] = True
        if temp_info:
            all_info[prop_number].update(temp_info)

    # All 3 POC (Tech POC, customer, end user)
    # are required for all but Phase I proposals
    # If any are missing for other Phases, annotate as "NOT FOUND"
    for k,v in all_info.items():
        for poc_header in poc_headers:
            if poc_header not in v:
                if v['Phase'] != "I":
                    v[poc_header]="NOT FOUND"
                    pass
                else:
                    v[poc_header]="N/A"
        # I want all the POC info adjecent in the final dataframe/csv, so
        # I make the change below so they're adjecent when sorted
        new_name = "Primary Technical Points of Contact (TPOCs)"
        v[new_name] = v.pop("Technical Points of Contact (TPOCs)")


    end = time.time()
    print(f"{len(target_files)} files in {end-start} seconds")

    results = pd.DataFrame.from_dict(all_info, orient="index")
    if not any(results):
        logger.critical("Something went wrong, results is empty")

    sorted_cols = results.columns.tolist()
    sorted_cols.sort(key=natural_keys)
    logger.debug("Sorted columns:")
    logger.debug(sorted_cols)
    results = results[['Type', 'Phase', 'Page Count']+ [c for c in sorted_cols if c not in ['Type', 'Phase', 'Page Count']]]
    # Sort columns in dataframe by name (alphanumerically)
    # results = results[sorted_cols]
    results.index.name = "Proposal ID"
    # Print and save resulting table

    pprint.pprint(results.T)
    results.to_csv(args.out)
    return len(results)
# ==============================================================================
if __name__ == "__main__":

    # Create the parser and add arguments
    parser = argparse.ArgumentParser(
        formatter_class=argparse.ArgumentDefaultsHelpFormatter)

    parser.add_argument('--budget-file',
                        '-b',
                        type=str,
                        default="budget",
                        help="Will parse files with this keyword as budget documents"
                        )

    parser.add_argument('--questions-file',
                        '-q',
                        type=str,
                        default="all_forms",
                        action="append",
                        help="Will parse files with this keyword for duration and firm+proposal questions"
                        )

    parser.add_argument('--sig-file',
                        '-s',
                        type=str,
                        default="mou",
                        action="append",
                        help="Will parse files with this keyword for POC signatures"
                        )

    parser.add_argument('--vol2-file',
                        type=str,
                        default="Vol2-proposal",
                        action="append",
                        help="Will count pages/slides in files containing this keyword"
                        )

    parser.add_argument('--file',
                        '-f',
                        action='append',
                        type=is_filename,
                        help="Files to parse.  Must be pdf/ppt/pptx")

    parser.add_argument('--directory',
                        '-d',
                        action='append',
                        type=is_directory,
                        help="Directories to search for files to parse")

    parser.add_argument('--keyword',
                        '-k',
                        action='append',
                        type=lower_str,
                        help="Parse only filenames containing ALL these keywords"
                        )

    parser.add_argument('--ocr',
                        '-o',
                        type=str2bool,
                        default=False,
                        help="Use OCR (Slower, but can parse scanned PDFs)"
                        )

    parser.add_argument('--out',
                        type=str,
                        default="proposals.csv",
                        help="Save results to this file (will be a csv.)"
                        )
    parser.add_argument('--max-value',
                        '-m',
                        type=int,
                        default=1250000,
                        help="Max dollar value for proposals.  Will print warning if value is exceeded"
                        )
    parser.add_argument('--log-level',
                        '-l',
                        type=int,
                        default=50,
                        help="log level (0-50) 50 is the default"
                        )
    args = parser.parse_args()

    if not args.file and not args.directory:
        print("Must provide at least one pdf file or directory (will recurse over all directory files.)")
        sys.exit(1)

    # Invocation correct; now load modules
    # Otherwise you force the user to wait for them to load,
    # then tell them the invocation is incorrect, what a waste of time.

    print("Invocation correct, loading standard modules")
    import time
    start = time.time()
    import pprint
    import re
    import math
    import logging

    print("Loading pandas and collections modules")
    import pandas as pd
    from collections import defaultdict

    # Note: the PPT module is lazily loaded in main()

    print("Loading fitz (pdf module)")
    import fitz
    if args.ocr:
        print("Loading OCR modules...")
        from pdf2image import convert_from_path
        from PIL import Image
        import pytesseract

    print("Done loading modules")

    # Configure the logger
    logging.basicConfig(level=args.log_level, format='%(message)s')
    logger = logging.getLogger(__name__)
    logger.setLevel(args.log_level)

    # Pre-compile re expressions
    safety_heading = re.compile("safety.*related.*deliverables")
    some_digits = re.compile('\d{5,}')
    four_digits = re.compile(r"(\d{4})")

    pd.set_option('display.width', 80)
    pd.set_option('display.max_colwidth', 80)
    # If you don't have tesseract executable in your PATH, include the following:
    #pytesseract.pytesseract.tesseract_cmd = r'C:\\Program Files\\Tesseract-OCR\\tesseract'
    pprint.pprint(args)

    ppt_extensions = ["ppt", "pptx"]
    valid_extensions = ppt_extensions + ["pdf"]

    # args is global so no need to pass it
    number_props = main()
    end = time.time()

    print(f"{number_props} proposals in {round(end-start,2)} seconds")