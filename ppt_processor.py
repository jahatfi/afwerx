import os
import pprint
import re
import shutil
import sys
import time
from pptx import Presentation

# TODO Update
sections = [
    "Method",
    "Approach"
]
#===============================================================================
def is_filename(filename:str):
    try:
        f = open(filename)
    except FileNotFoundError:
        raise argparse.ArgumentTypeError(f"Cannot open file'{filename}'")
    finally:
        try:
            f.close()
        except Exception:
            pass
    return filename
#===============================================================================
def is_directory(dir:str):
    if not os.path.isdir(dir):
        raise argparse.ArgumentTypeError(f"'{dir}' is not a valid directory")
    return os.path.abspath(dir)
# ==============================================================================
def main(args):
    target_files = set()
    dirs = []
    valid_extensions = ["ppt", "pptx"]
    if not args.directory:
        args.directory = set()
    if not args.file:
        args.file = set()
    for dir in args.directory:
        dirs.append(dir)

    for each_file in args.file:
        file_extension = each_file.split(".")[-1]
        if file_extension in valid_extensions:
            target_files.add(each_file)
        else:
            print(f"Skipping {each_file} with extension {file_extension}")

    # Count the total number of files to be parsed.
    total_files = len(target_files)
    for source_dir in dirs:
        for (root, _, files) in os.walk(source_dir):
            for file_name in files:
                if file_extension in valid_extensions:
                    target_files.add(each_file)  

    print(f"Parsing {len(target_files)} files. This could take a few seconds.")
    # Reset the counter.  It will be incremented as each file is parsed.
    total_files = 0

    # Here is the serial (non-parallel) approach.  Slow, but it works.
    start = time.time()

    for filename in sorted(target_files):
        print(f"{filename}".center(80, '-'))
        prs = Presentation(filename)
        for slide in prs.slides:
            try:
                title = slide.shapes.title.text
            except AttributeError:
                title = slide.shapes[0].text
            print(title)
            # TODO: lowercase and compare to remaining list
    end = time.time()
    print(f"{total_files} files in {end-start} seconds")
# ==============================================================================
# ==============================================================================
if __name__ == "__main__":
    import argparse

    # Create the parser and add arguments
    parser = argparse.ArgumentParser()
    failed = False
    # More libraries are loaded if invocation is correct

    # Add an optional argument for the output file,
    # open in 'write' mode and and specify encoding
    parser.add_argument('--output_dir', 
                        '-o', 
                        type=is_directory, 
                        default=".", 
                        help="Name of directory to store results.")

    parser.add_argument('--file',
                        '-f',
                        action='append', 
                        type=is_filename)

    parser.add_argument('--directory',
                        '-d', 
                        action='append', 
                        type=is_directory)                           
    args = parser.parse_args()
    pprint.pprint(args)


    if not args.file and not args.directory:
        print("Must provide at least one pdf file or directory (will recurse over all directory files.)")
        sys.exit(1)

    original_dir = os.getcwd()
    slide = main(args)